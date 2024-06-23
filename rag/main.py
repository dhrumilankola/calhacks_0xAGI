# rag_tutor_with_optimized_vector_store.py

import os
import nbformat
import pickle
from nbformat import read
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed
from langchain_community.document_loaders import WebBaseLoader, PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_core.output_parsers import StrOutputParser
from langchain.prompts import ChatPromptTemplate
from langchain_community.chat_models import ChatOllama
from langchain_core.runnables import RunnableLambda, RunnablePassthrough
from langchain.schema import Document

DOCUMENTS_PATH = 'documents.pkl'
EMBEDDINGS_PATH = 'embeddings.pkl'

# Function to extract text from PDF files
def extract_text_from_pdf(pdf_path):
    pdf_loader = PyPDFLoader(pdf_path)
    return [Document(page_content=doc.page_content, metadata={"source": pdf_path}) for doc in pdf_loader.load()]

# Function to extract text from PPTX files
def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    pptx_text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pptx_text += shape.text + "\n\n"
    return [Document(page_content=pptx_text, metadata={"source": pptx_path})]

# Function to extract text from Jupyter Notebooks
def extract_text_from_notebook(notebook_path):
    with open(notebook_path, 'r', encoding='utf-8') as f:
        notebook = read(f, as_version=4)
    notebook_text = ""
    for cell in notebook.cells:
        if cell.cell_type == 'markdown' or cell.cell_type == 'code':
            notebook_text += cell.source + "\n\n"
    return [Document(page_content=notebook_text, metadata={"source": notebook_path})]

# Function to load and preprocess documents
def load_and_preprocess_documents():
    data = []

    # Extract text from dataset folders
    pdf_folder = 'dataset/books'
    pptx_folder = 'dataset/slides'
    notebook_folder = 'dataset/code demo'

    with ThreadPoolExecutor() as executor:
        futures = []
        for pdf_file in os.listdir(pdf_folder):
            if pdf_file.endswith(".pdf"):
                pdf_path = os.path.join(pdf_folder, pdf_file)
                futures.append(executor.submit(extract_text_from_pdf, pdf_path))

        for pptx_file in os.listdir(pptx_folder):
            if pptx_file.endswith(".pptx"):
                pptx_path = os.path.join(pptx_folder, pptx_file)
                futures.append(executor.submit(extract_text_from_pptx, pptx_path))

        for notebook_file in os.listdir(notebook_folder):
            if notebook_file.endswith(".ipynb"):
                notebook_path = os.path.join(notebook_folder, notebook_file)
                futures.append(executor.submit(extract_text_from_notebook, notebook_path))

        for future in as_completed(futures):
            data.extend(future.result())

    # Split documents into smaller chunks
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=0)
    all_splits = text_splitter.split_documents(data)

    return all_splits

# Function to initialize the vector store and retriever
def initialize_vector_store(documents, embeddings):
    # Store documents in ChromaDB vector store
    vectorstore = Chroma.from_documents(
        documents=documents,
        collection_name="rag-chroma",
        embedding=embeddings,
    )

    return vectorstore

# Save documents and embeddings to files
def save_documents_and_embeddings(documents, embeddings, doc_path, emb_path):
    with open(doc_path, 'wb') as doc_file:
        pickle.dump(documents, doc_file)
    with open(emb_path, 'wb') as emb_file:
        pickle.dump(embeddings, emb_file)

# Load documents and embeddings from files
def load_documents_and_embeddings(doc_path, emb_path):
    with open(doc_path, 'rb') as doc_file:
        documents = pickle.load(doc_file)
    with open(emb_path, 'rb') as emb_file:
        embeddings = pickle.load(emb_file)
    return documents, embeddings

# Function to set up the RAG chain
def setup_rag_chain(retriever):
    # Robust prompt template
    template = """
    You are an AI tutor designed to assist students by providing accurate, detailed, and helpful answers to their questions based on the provided context. Use the information from the following documents to formulate your response.

    Context:
    {context}

    Question:
    {question}

    Answer the question as thoroughly as possible, including explanations, examples, and any relevant information that will help the student understand the topic better. If the context does not provide enough information, acknowledge this and provide a logical and educated guess based on your general knowledge.
    """
    prompt = ChatPromptTemplate.from_template(template)

    # Local LLM configuration
    ollama_llm = "phi3:3.8b"
    model_local = ChatOllama(model=ollama_llm)

    # Chain setup
    chain = (
        {"context": retriever, "question": RunnablePassthrough()}
        | prompt
        | model_local
        | StrOutputParser()
    )

    return chain

# Main function to handle queries
def handle_query(question):
    # Check if documents and embeddings already exist
    if os.path.exists(DOCUMENTS_PATH) and os.path.exists(EMBEDDINGS_PATH):
        documents, embeddings = load_documents_and_embeddings(DOCUMENTS_PATH, EMBEDDINGS_PATH)
    else:
        # Load and preprocess documents
        documents = load_and_preprocess_documents()
        
        # Generate embeddings
        embeddings = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")

        # Save documents and embeddings for future use
        save_documents_and_embeddings(documents, embeddings, DOCUMENTS_PATH, EMBEDDINGS_PATH)

    vectorstore = initialize_vector_store(documents, embeddings)
    retriever = vectorstore.as_retriever()
    
    # Setup RAG chain
    chain = setup_rag_chain(retriever)

    # Invoke chain with the question
    response = chain.invoke(question)

    return response


# Example usage
if __name__ == "__main__":
    question = "Explain about transformers, also explain me the transformers notebook code demo provided by the professor"
    response = handle_query(question)
    print(response)

